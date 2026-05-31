#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT自动重新设计工具
使用AI帮助优化PPT的布局、设计和内容排列
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import sys
from pathlib import Path

class PPTRedesigner:
    def __init__(self, ppt_file):
        """初始化PPT重新设计工具"""
        self.ppt_file = ppt_file
        self.prs = None
        self.load_ppt()
    
    def load_ppt(self):
        """加载PPT文件"""
        try:
            self.prs = Presentation(self.ppt_file)
            print(f"✅ 成功加载PPT文件: {self.ppt_file}")
            print(f"📊 总共有 {len(self.prs.slides)} 张幻灯片")
        except Exception as e:
            print(f"❌ 加载PPT失败: {e}")
            sys.exit(1)
    
    def analyze_slides(self):
        """分析PPT中的所有幻灯片内容"""
        slides_info = []
        for idx, slide in enumerate(self.prs.slides, 1):
            slide_data = {
                'slide_number': idx,
                'shapes': [],
                'text_content': [],
                'images': []
            }
            
            for shape in slide.shapes:
                # 提取文本
                if hasattr(shape, "text") and shape.text.strip():
                    slide_data['text_content'].append(shape.text)
                    slide_data['shapes'].append({
                        'type': 'text',
                        'content': shape.text,
                        'position': (shape.left, shape.top),
                        'size': (shape.width, shape.height)
                    })
                
                # 提取图片信息
                if shape.shape_type == 13:  # 图片类型
                    slide_data['images'].append({
                        'type': 'image',
                        'position': (shape.left, shape.top),
                        'size': (shape.width, shape.height)
                    })
            
            slides_info.append(slide_data)
        
        return slides_info
    
    def print_analysis(self):
        """打印分析结果"""
        slides_info = self.analyze_slides()
        print("\n" + "="*60)
        print("📋 PPT内容分析")
        print("="*60)
        
        for slide_data in slides_info:
            print(f"\n📄 第 {slide_data['slide_number']} 张幻灯片:")
            
            if slide_data['text_content']:
                print("  📝 文本内容:")
                for text in slide_data['text_content']:
                    print(f"     - {text[:50]}{'...' if len(text) > 50 else ''}")
            
            if slide_data['images']:
                print(f"  🖼️  图片: {len(slide_data['images'])} 张")
    
    def apply_theme(self, theme='modern'):
        """
        应用设计主题
        theme: 'modern' (现代), 'professional' (专业), 'creative' (创意)
        """
        print(f"\n🎨 应用 '{theme}' 主题...")
        
        # 定义主题色
        themes = {
            'modern': {
                'primary': RGBColor(0, 120, 215),      # 蓝色
                'secondary': RGBColor(50, 50, 50),     # 深灰
                'accent': RGBColor(255, 150, 0)        # 橙色
            },
            'professional': {
                'primary': RGBColor(31, 78, 121),      # 深蓝
                'secondary': RGBColor(89, 89, 89),     # 灰色
                'accent': RGBColor(192, 0, 0)          # 红色
            },
            'creative': {
                'primary': RGBColor(162, 109, 173),    # 紫色
                'secondary': RGBColor(100, 149, 237),  # 矢车菊蓝
                'accent': RGBColor(255, 165, 0)        # 金色
            }
        }
        
        color_scheme = themes.get(theme, themes['modern'])
        
        for slide_idx, slide in enumerate(self.prs.slides):
            for shape in slide.shapes:
                # 调整文本颜色
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        # 标题使用主色
                        if slide_idx == 0 or shape.top < Inches(1.5):
                            for run in paragraph.runs:
                                run.font.color.rgb = color_scheme['primary']
                        # 正文使用次色
                        else:
                            for run in paragraph.runs:
                                run.font.color.rgb = color_scheme['secondary']
        
        print("✅ 主题应用完成")
    
    def optimize_layout(self):
        """优化幻灯片布局"""
        print("\n📐 优化布局...")
        
        for slide in self.prs.slides:
            text_shapes = [s for s in slide.shapes if hasattr(s, "text_frame")]
            image_shapes = [s for s in slide.shapes if s.shape_type == 13]
            
            # 如果有图片和文本，调整它们的相对位置
            if text_shapes and image_shapes:
                # 文本放左边，图片放右边
                for text_shape in text_shapes:
                    text_shape.left = Inches(0.5)
                    text_shape.width = Inches(4.5)
                
                for img_shape in image_shapes:
                    img_shape.left = Inches(5.5)
                    img_shape.width = Inches(4)
        
        print("✅ 布局优化完成")
    
    def save_redesigned_ppt(self, output_file=None):
        """保存重新设计后的PPT"""
        if output_file is None:
            base_name = os.path.splitext(self.ppt_file)[0]
            output_file = f"{base_name}_redesigned.pptx"
        
        try:
            self.prs.save(output_file)
            print(f"\n✅ 重新设计的PPT已保存: {output_file}")
            return output_file
        except Exception as e:
            print(f"❌ 保存失败: {e}")
            return None
    
    def redesign(self, theme='modern', optimize=True):
        """执行完整的PPT重新设计"""
        print("\n" + "="*60)
        print("🚀 开始PPT重新设计流程")
        print("="*60)
        
        # 分析内容
        self.print_analysis()
        
        # 应用主题
        self.apply_theme(theme)
        
        # 优化布局
        if optimize:
            self.optimize_layout()
        
        # 保存
        output = self.save_redesigned_ppt()
        
        print("\n" + "="*60)
        print("✨ PPT重新设计完成！")
        print("="*60)
        
        return output


def main():
    """主函数"""
    # 查找PPT文件
    ppt_files = list(Path('.').glob('*.pptx')) + list(Path('.').glob('**/*.pptx'))
    
    if not ppt_files:
        print("❌ 未找到PPT文件，请确保.pptx文件在当前目录或子目录中")
        sys.exit(1)
    
    print(f"🔍 找到 {len(ppt_files)} 个PPT文件")
    
    for ppt_file in ppt_files:
        print(f"\n处理: {ppt_file}")
        
        # 创建重新设计器
        redesigner = PPTRedesigner(str(ppt_file))
        
        # 执行重新设计（可选择不同主题）
        # 可用主题: 'modern' (现代), 'professional' (专业), 'creative' (创意)
        redesigner.redesign(theme='modern', optimize=True)


if __name__ == '__main__':
    main()